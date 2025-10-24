from pptx import Presentation
import os
import json
import time
import openai
from openai import OpenAI

# OpenAI API配置
# 请在这里设置您的API key
OPENAI_API_KEY = ""  # 请替换为您的实际API key
client = OpenAI(api_key=OPENAI_API_KEY)

def extract_ppt_content(ppt_path, output_dir="ppt_output", save_images=True):
    """
    提取PPTX文件的纯文本内容，按页数一一对应
    
    Args:
        ppt_path (str): PPT文件路径
        output_dir (str): 输出目录
        save_images (bool): 是否保存图片文件
    
    Returns:
        list: 包含所有幻灯片信息的列表
    """
    # 检查文件是否存在
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")
    
    prs = Presentation(ppt_path)
    os.makedirs(output_dir, exist_ok=True)
    
    all_slides = []
    
    for i, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "slide_number": i,
            "texts": [],
            "images": []
        }
        
        for shape in slide.shapes:
            # 提取纯文字内容
            if hasattr(shape, "text") and shape.text.strip():
                slide_info["texts"].append({
                    "content": shape.text.strip()
                })
            
            # 提取图片（如果需要）
            if shape.shape_type == 13:  # 13表示图片类型
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_name = f"slide_{i}_img_{len(slide_info['images']) + 1}.{image_ext}"
                    image_path = os.path.join(output_dir, image_name)
                    
                    if save_images:
                        with open(image_path, "wb") as f:
                            f.write(image_bytes)
                    
                    slide_info["images"].append({
                        "filename": image_name
                    })
                except Exception as e:
                    print(f"⚠️ 提取第{i}页图片时出错: {e}")
        
        all_slides.append(slide_info)
    
    # 导出为JSON
    json_path = os.path.join(output_dir, "ppt_content.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(all_slides, f, ensure_ascii=False, indent=2)
    
    print(f"✅ 提取完成，共 {len(all_slides)} 页。结果保存在：{json_path}")
    
    # 打印摘要信息
    total_texts = sum(len(slide["texts"]) for slide in all_slides)
    total_images = sum(len(slide["images"]) for slide in all_slides)
    print(f"📊 统计信息: 文本块 {total_texts} 个，图片 {total_images} 张")
    
    return all_slides

def translate_to_japanese(text):
    """
    使用OpenAI API将文本翻译成日语
    
    Args:
        text (str): 要翻译的文本
    
    Returns:
        str: 翻译后的日语文本
    """
    try:
        # 使用OpenAI API进行翻译
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",  # 或者使用 "gpt-4" 获得更好的翻译质量
            messages=[
                {
                    "role": "system",
                    "content": "你是一个专业的中文到日语翻译助手。请将用户提供的中文文本准确翻译成日语，保持原文的语气和含义。对于专业术语，请使用准确的日语表达。"
                },
                {
                    "role": "user",
                    "content": f"请将以下中文文本翻译成日语：\n{text}"
                }
            ],
            temperature=0.3,  # 较低的temperature确保翻译的一致性
            max_tokens=1000
        )
        
        translated_text = response.choices[0].message.content.strip()
        return translated_text
        
    except Exception as e:
        print(f"⚠️ OpenAI翻译失败: {e}")
        return text  # 翻译失败时返回原文

def batch_translate_slides(slides_data):
    """
    批量翻译幻灯片内容
    
    Args:
        slides_data (list): 幻灯片数据
    
    Returns:
        list: 翻译后的幻灯片数据
    """
    translated_slides = []
    
    for slide in slides_data:
        translated_slide = {
            "slide_number": slide["slide_number"],
            "texts": [],
            "images": slide["images"]  # 图片信息保持不变
        }
        
        for text_item in slide["texts"]:
            original_text = text_item["content"]
            translated_text = translate_to_japanese(original_text)
            
            translated_slide["texts"].append({
                "content": translated_text,
                "original_content": original_text  # 保留原文
            })
            
            # 添加延迟避免API限制
            time.sleep(1)  # OpenAI API建议的延迟时间
        
        translated_slides.append(translated_slide)
        print(f"✅ 第{slide['slide_number']}页翻译完成")
    
    return translated_slides

def replace_ppt_text_with_translation(ppt_path, translated_data, output_path):
    """
    将PPT中的文本替换为翻译后的日语文本
    
    Args:
        ppt_path (str): 原PPT文件路径
        translated_data (list): 翻译后的数据
        output_path (str): 输出PPT文件路径
    """
    prs = Presentation(ppt_path)
    
    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        translated_slide = next((s for s in translated_data if s["slide_number"] == slide_number), None)
        
        if not translated_slide:
            continue
        
        text_index = 0
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if text_index < len(translated_slide["texts"]):
                    # 替换为翻译后的文本
                    shape.text = translated_slide["texts"][text_index]["content"]
                    text_index += 1
    
    # 保存新的PPT文件
    prs.save(output_path)
    print(f"✅ 翻译后的PPT已保存到: {output_path}")

def print_content_summary(slides_data):
    """
    打印内容摘要
    """
    print("\n📋 内容摘要:")
    for slide in slides_data:
        slide_num = slide["slide_number"]
        text_count = len(slide["texts"])
        image_count = len(slide["images"])
        print(f"第{slide_num}页: {text_count}个文本块, {image_count}张图片")
        
        # 显示前50个字符的文本预览
        if slide["texts"]:
            preview = slide["texts"][0]["content"][:50] + "..." if len(slide["texts"][0]["content"]) > 50 else slide["texts"][0]["content"]
            print(f"  文本预览: {preview}")

# 示例用法
if __name__ == "__main__":
    # 从trip7_ppt_translation/chinese目录读取PPT文件
    chinese_dir = "trip7_ppt_translation/chinese"
    ppt_files = [f for f in os.listdir(chinese_dir) if f.endswith('.pptx')]
    
    if not ppt_files:
        print("❌ 在chinese目录中未找到PPTX文件")
        exit(1)
    
    # 处理第一个找到的PPTX文件
    ppt_file = os.path.join(chinese_dir, ppt_files[0])
    print(f"📂 处理文件: {ppt_file}")
    
    try:
        # 步骤1: 提取内容
        print("🔍 步骤1: 提取PPT内容...")
        slides_data = extract_ppt_content(ppt_file, output_dir="trip7_ppt_translation/extracted_content")
        
        # 步骤2: 翻译内容
        print("🌐 步骤2: 翻译内容到日语...")
        translated_data = batch_translate_slides(slides_data)
        
        # 保存翻译后的JSON
        translated_json_path = "trip7_ppt_translation/extracted_content/translated_content.json"
        with open(translated_json_path, "w", encoding="utf-8") as f:
            json.dump(translated_data, f, ensure_ascii=False, indent=2)
        print(f"💾 翻译结果已保存到: {translated_json_path}")
        
        # 步骤3: 创建日语版PPT
        print("📝 步骤3: 创建日语版PPT...")
        japanese_dir = "trip7_ppt_translation/japanese"
        os.makedirs(japanese_dir, exist_ok=True)
        
        output_ppt_path = os.path.join(japanese_dir, f"japanese_{ppt_files[0]}")
        replace_ppt_text_with_translation(ppt_file, translated_data, output_ppt_path)
        
        print("\n🎉 翻译工作流完成！")
        print(f"📊 处理了 {len(slides_data)} 页幻灯片")
        print(f"📁 日语版PPT: {output_ppt_path}")
        
    except FileNotFoundError as e:
        print(f"❌ 错误: {e}")
        print("请确保PPT文件存在于chinese目录中")
    except Exception as e:
        print(f"❌ 处理过程中出现错误: {e}")