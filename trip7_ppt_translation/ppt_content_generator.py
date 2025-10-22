from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import json
import os
import shutil
from PIL import Image

def create_ppt_from_json(json_path, output_ppt_path, template_ppt_path=None):
    """
    从日语JSON文件生成PPT文件，按页数对应模板页面
    
    Args:
        json_path (str): 日语JSON文件路径
        output_ppt_path (str): 输出PPT文件路径
        template_ppt_path (str): 模板PPT文件路径（可选）
    """
    try:
        # 读取JSON数据
        print("📖 读取日语JSON文件...")
        with open(json_path, 'r', encoding='utf-8') as f:
            slides_data = json.load(f)
        
        # 创建新的演示文稿或使用模板
        if template_ppt_path and os.path.exists(template_ppt_path):
            print(f"📋 使用模板文件: {template_ppt_path}")
            
            # 直接复制模板文件作为基础，保持完整格式
            shutil.copy2(template_ppt_path, output_ppt_path)
            prs = Presentation(output_ppt_path)
            
            print(f"🔄 开始修改 {len(slides_data)} 页幻灯片...")
            
            # 确保幻灯片数量匹配
            while len(prs.slides) < len(slides_data):
                # 复制最后一页作为新页面
                last_slide = prs.slides[-1]
                slide_layout = last_slide.slide_layout
                new_slide = prs.slides.add_slide(slide_layout)
            
            # 删除多余的幻灯片（如果有）
            while len(prs.slides) > len(slides_data):
                rId = prs.slides._sldIdLst[-1].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[-1]
            
            # 替换每页的文本内容
            for i, slide_data in enumerate(slides_data):
                slide_number = slide_data.get('slide_number', i+1)
                print(f"📄 修改第 {slide_number} 页...")
                
                slide = prs.slides[i]
                texts = slide_data.get('texts', [])
                if texts:
                    replace_slide_text_safe(slide, texts)
        else:
            print("📋 创建新的演示文稿")
            prs = Presentation()
            # 删除默认的空白幻灯片
            if len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
            
            print(f"🔄 开始生成 {len(slides_data)} 页幻灯片...")
            
            # 为每个幻灯片数据创建幻灯片
            for i, slide_data in enumerate(slides_data):
                slide_number = slide_data.get('slide_number', i+1)
                print(f"📄 生成第 {slide_number} 页...")
                
                # 添加空白幻灯片
                blank_slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 添加文本
                texts = slide_data.get('texts', [])
                if texts:
                    create_text_boxes(slide, texts)
        
        # 保存PPT文件
        print(f"💾 保存PPT文件: {output_ppt_path}")
        prs.save(output_ppt_path)
        
        print(f"✅ 日语PPT文件生成完成！")
        print(f"📂 输入文件: {json_path}")
        print(f"📋 输出文件: {output_ppt_path}")
        print(f"📊 共生成 {len(slides_data)} 页幻灯片")
        
    except FileNotFoundError as e:
        print(f"❌ 文件未找到: {e}")
    except json.JSONDecodeError:
        print(f"❌ JSON文件格式错误")
    except Exception as e:
        print(f"❌ 生成PPT过程中出现错误: {e}")

def replace_slide_text_safe(slide, texts):
    """
    安全地替换幻灯片中的文本内容
    """
    try:
        text_shapes = []
        
        # 收集所有可编辑的文本形状
        for shape in slide.shapes:
            if (hasattr(shape, 'text_frame') and 
                shape.text_frame and 
                hasattr(shape.text_frame, 'paragraphs')):
                text_shapes.append(shape)
        
        # 按位置排序（从上到下，从左到右）
        text_shapes.sort(key=lambda x: (getattr(x, 'top', 0), getattr(x, 'left', 0)))
        
        # 替换文本内容
        for i, text_info in enumerate(texts):
            content = text_info.get('content', '').strip()
            if not content:
                continue
            
            if i < len(text_shapes):
                shape = text_shapes[i]
                
                # 保存原有格式
                original_font_name = None
                original_font_size = None
                
                try:
                    if shape.text_frame.paragraphs:
                        first_para = shape.text_frame.paragraphs[0]
                        if first_para.runs:
                            first_run = first_para.runs[0]
                            original_font_name = first_run.font.name
                            original_font_size = first_run.font.size
                except:
                    pass
                
                # 替换文本内容
                shape.text = content
                
                # 重新应用格式
                try:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 优先使用日语字体，但保持原有大小
                            if original_font_name and 'Arial' not in str(original_font_name):
                                run.font.name = original_font_name
                            else:
                                # 使用更通用的日语字体
                                run.font.name = 'Microsoft YaHei UI'
                            
                            if original_font_size:
                                run.font.size = original_font_size
                except Exception as font_error:
                    print(f"⚠️  设置字体时出错: {font_error}")
            else:
                # 如果文本数量超过现有文本框，创建新的文本框
                create_additional_textbox(slide, content, i)
    
    except Exception as e:
        print(f"⚠️  替换文本时出错: {e}")

def create_additional_textbox(slide, content, index):
    """
    创建额外的文本框
    """
    try:
        # 计算位置
        left = Inches(1)
        top = Inches(1.5 + index * 0.8)
        width = Inches(8)
        height = Inches(0.6)
        
        # 添加文本框
        textbox = slide.shapes.add_textbox(
            left=left,
            top=top,
            width=width,
            height=height
        )
        
        textbox.text = content
        
        # 设置字体
        for paragraph in textbox.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Microsoft YaHei UI'
                run.font.size = Pt(16)
    
    except Exception as e:
        print(f"⚠️  创建额外文本框时出错: {e}")

def create_text_boxes(slide, texts):
    """
    创建文本框（无模板时使用）
    """
    try:
        # 计算文本框位置
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(0.8)
        
        for i, text_info in enumerate(texts):
            content = text_info.get('content', '').strip()
            if not content:
                continue
            
            # 计算当前文本框位置
            current_top = top + (height * i)
            
            # 添加文本框
            textbox = slide.shapes.add_textbox(
                left=left,
                top=current_top,
                width=width,
                height=height
            )
            
            text_frame = textbox.text_frame
            text_frame.text = content
            
            # 格式化文本
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Microsoft YaHei UI'
                    run.font.size = Pt(16)
                    
    except Exception as e:
        print(f"⚠️  创建文本框时出错: {e}")

def validate_generated_ppt(ppt_path):
    """
    验证生成的PPT文件
    """
    try:
        prs = Presentation(ppt_path)
        
        print(f"\n🔍 验证生成的PPT文件: {ppt_path}")
        print(f"📊 总幻灯片数: {len(prs.slides)}")
        print(f"📐 幻灯片尺寸: {prs.slide_width} x {prs.slide_height}")
        
        # 检查每页内容
        total_shapes = 0
        for i, slide in enumerate(prs.slides, 1):
            shape_count = len(slide.shapes)
            total_shapes += shape_count
            print(f"第{i}页: {shape_count}个形状")
        
        print(f"📈 总形状数: {total_shapes}")
        print("✅ PPT文件验证完成")
        
    except Exception as e:
        print(f"❌ 验证PPT文件时出错: {e}")

if __name__ == "__main__":
    # 设置文件路径
    japanese_json = "trip7_ppt_translation/extracted_content/ppt_content.japanese.json"
    output_ppt = "trip7_ppt_translation/japanese/Web3与元宇宙_日语版.pptx"
    template_ppt = "trip7_ppt_translation/chinese/第1课-课程简介-Web3简介(2课时)-20251020.pptx"  # 模板文件路径
    
    print("🚀 开始生成日语PPT文件...")
    print("=" * 50)
    
    # 检查输入文件是否存在
    if not os.path.exists(japanese_json):
        print(f"❌ 日语JSON文件不存在: {japanese_json}")
        print("请先运行ppt_content_filler.py生成日语JSON文件")
        exit(1)
    
    # 创建输出目录
    os.makedirs(os.path.dirname(output_ppt), exist_ok=True)
    
    # 检查模板文件
    if os.path.exists(template_ppt):
        print(f"📋 找到模板文件: {template_ppt}")
        create_ppt_from_json(japanese_json, output_ppt, template_ppt)
    else:
        print(f"⚠️  模板文件不存在: {template_ppt}")
        print("将使用默认格式生成PPT")
        create_ppt_from_json(japanese_json, output_ppt)
    
    # 验证生成的文件
    if os.path.exists(output_ppt):
        validate_generated_ppt(output_ppt)
    
    print("\n🎉 日语PPT生成任务完成！")