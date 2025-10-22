from pptx import Presentation
from PIL import Image
import io
import os
import json
from pptx.util import Inches

def extract_ppt_content(ppt_path, output_dir="ppt_output", save_images=True):
    """
    提取PPTX文件的内容（文字和图片），按页数一一对应，包含样式和位置信息
    
    Args:
        ppt_path (str): PPT文件路径
        output_dir (str): 输出目录
        save_images (bool): 是否保存图片文件
    
    Returns:
        list: 包含所有幻灯片信息的列表
    """
    # 检查文件是否存在
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")
    
    prs = Presentation(ppt_path)
    os.makedirs(output_dir, exist_ok=True)
    
    all_slides = []
    
    for i, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "slide_number": i,
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "texts": [],
            "images": []
        }
        
        for shape in slide.shapes:
            # 提取文字及其样式和位置信息
            if hasattr(shape, "text") and shape.text.strip():
                text_info = {
                    "content": shape.text.strip(),
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                }
                
                # 安全获取z_order
                try:
                    text_info["z_order"] = getattr(shape.element, 'z_order', 0) if hasattr(shape, 'element') else 0
                except:
                    text_info["z_order"] = 0
                
                # 提取文字样式信息
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text_info["style"] = {
                        "paragraphs": []
                    }
                    
                    for paragraph in shape.text_frame.paragraphs:
                        para_info = {
                            "text": paragraph.text,
                            "alignment": str(paragraph.alignment) if paragraph.alignment else None,
                            "runs": []
                        }
                        
                        for run in paragraph.runs:
                            run_info = {
                                "text": run.text,
                                "font_name": run.font.name if run.font.name else None,
                                "font_size": run.font.size.pt if run.font.size else None,
                                "bold": run.font.bold,
                                "italic": run.font.italic,
                                "underline": str(run.font.underline) if run.font.underline else None
                            }
                            
                            # 安全提取字体颜色
                            try:
                                if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                    run_info["color"] = str(run.font.color.rgb)
                                else:
                                    run_info["color"] = None
                            except Exception:
                                run_info["color"] = None
                            
                            para_info["runs"].append(run_info)
                        
                        text_info["style"]["paragraphs"].append(para_info)
                
                slide_info["texts"].append(text_info)
            
            # 提取图片及其位置信息
            if shape.shape_type == 13:  # 13表示图片类型
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_name = f"slide_{i}_img_{len(slide_info['images']) + 1}.{image_ext}"
                    image_path = os.path.join(output_dir, image_name)
                    
                    if save_images:
                        with open(image_path, "wb") as f:
                            f.write(image_bytes)
                    
                    image_info = {
                        "filename": image_name,
                        "position": {
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height
                        }
                    }
                    
                    # 安全获取z_order
                    try:
                        image_info["z_order"] = getattr(shape.element, 'z_order', 0) if hasattr(shape, 'element') else 0
                    except:
                        image_info["z_order"] = 0
                    
                    # 安全获取原始尺寸
                    try:
                        image_info["original_size"] = {
                            "width": image.size[0] if hasattr(image, 'size') else None,
                            "height": image.size[1] if hasattr(image, 'size') else None
                        }
                    except:
                        image_info["original_size"] = {"width": None, "height": None}
                    
                    slide_info["images"].append(image_info)
                except Exception as e:
                    print(f"⚠️ 提取第{i}页图片时出错: {e}")
        
        # 按z_order排序，保持原有的层次关系
        slide_info["texts"].sort(key=lambda x: x.get("z_order", 0))
        slide_info["images"].sort(key=lambda x: x.get("z_order", 0))
        
        all_slides.append(slide_info)
    
    # 导出为JSON
    json_path = os.path.join(output_dir, "ppt_content.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(all_slides, f, ensure_ascii=False, indent=2)
    
    print(f"✅ 提取完成，共 {len(all_slides)} 页。结果保存在：{json_path}")
    
    # 打印摘要信息
    total_texts = sum(len(slide["texts"]) for slide in all_slides)
    total_images = sum(len(slide["images"]) for slide in all_slides)
    print(f"📊 统计信息: 文本块 {total_texts} 个，图片 {total_images} 张")
    
    return all_slides

def print_content_summary(slides_data):
    """
    打印内容摘要
    """
    print("\n📋 内容摘要:")
    for slide in slides_data:
        slide_num = slide["slide_number"]
        text_count = len(slide["texts"])
        image_count = len(slide["images"])
        print(f"第{slide_num}页: {text_count}个文本块, {image_count}张图片")
        
        # 显示前50个字符的文本预览
        if slide["texts"]:
            preview = slide["texts"][0]["content"][:50] + "..." if len(slide["texts"][0]["content"]) > 50 else slide["texts"][0]["content"]
            print(f"  文本预览: {preview}")
            
            # 显示位置信息
            pos = slide["texts"][0]["position"]
            print(f"  位置: left={pos['left']}, top={pos['top']}, width={pos['width']}, height={pos['height']}")

# 示例用法
if __name__ == "__main__":
    # 从trip7_ppt_translation/chinese目录读取PPT文件
    chinese_dir = "trip7_ppt_translation/chinese"  # 修改这里
    ppt_files = [f for f in os.listdir(chinese_dir) if f.endswith('.pptx')]
    
    if not ppt_files:
        print("❌ 在chinese目录中未找到PPTX文件")
        exit(1)
    
    # 处理第一个找到的PPTX文件
    ppt_file = os.path.join(chinese_dir, ppt_files[0])
    print(f"📂 处理文件: {ppt_file}")
    
    try:
        # 提取内容到trip7_ppt_translation/extracted_content目录
        slides_data = extract_ppt_content(ppt_file, output_dir="trip7_ppt_translation/extracted_content")  # 修改这里
        
        # 打印摘要
        print_content_summary(slides_data)
        
    except FileNotFoundError as e:
        print(f"❌ 错误: {e}")
        print("请确保PPT文件存在于chinese目录中")
    except Exception as e:
        print(f"❌ 处理过程中出现错误: {e}")